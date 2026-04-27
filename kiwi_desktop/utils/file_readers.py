"""Helpers for extracting text samples from supported file types."""

from __future__ import annotations

import re
from pathlib import Path

MAX_EXTRACT_BYTES = 50 * 1024 * 1024  # 50MB


def _read_pptx(path: Path, max_chars: int) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        chunks: list[str] = []
        total = 0
        for index, slide in enumerate(prs.slides, start=1):
            title = ""
            if slide.shapes.title is not None:
                title = (slide.shapes.title.text or "").strip()
            if not title:
                title = "Untitled"

            body_parts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = (shape.text or "").strip()
                    if text:
                        body_parts.append(text)
            body = "\n".join(body_parts).strip()

            notes = ""
            notes_slide = getattr(slide, "notes_slide", None)
            if notes_slide is not None and notes_slide.notes_text_frame is not None:
                notes = (notes_slide.notes_text_frame.text or "").strip()

            slide_chunk = f"## Slide {index}: {title}"
            if body:
                slide_chunk = f"{slide_chunk}\n{body}"
            if notes:
                slide_chunk = f"{slide_chunk}\n> Notes: {notes}"
            chunks.append(slide_chunk)

            total += len(slide_chunk)
            if total >= max_chars:
                break
        return "\n\n".join(chunks)[:max_chars]
    except ImportError:
        return ""
    except Exception:
        return ""


def _read_pdf(path: Path, max_chars: int) -> str:
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(str(path))
        chunks = []
        total = 0
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text").strip()
            if text:
                chunks.append(f"[Page {page_num + 1}]\n{text}")
                total += len(text)
            if total >= max_chars:
                break
        doc.close()
        return "\n\n".join(chunks)[:max_chars]
    except ImportError:
        # PyMuPDF not installed — fall back to regex
        return _read_pdf_regex_fallback(path, max_chars)
    except Exception:
        return ""


def _read_pdf_regex_fallback(path: Path, max_chars: int) -> str:
    # Kept as emergency fallback only
    try:
        raw = path.read_bytes()
        matches = re.findall(b"BT(.+?)ET", raw, re.DOTALL)
        text = b" ".join(matches).decode("latin-1", errors="ignore")
        text = re.sub(r"\(([^)]*)\)", r"\1 ", text)
        text = re.sub(r"[^\x20-\x7E\n]", " ", text)
        return text[:max_chars]
    except Exception:
        return ""


def _read_docx(path: Path, max_chars: int) -> str:
    try:
        from docx import Document

        doc = Document(str(path))
        chunks: list[str] = []

        # Extract paragraphs preserving heading levels.
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name.lower() if para.style else ""
            if "heading 1" in style:
                chunks.append(f"# {text}")
            elif "heading 2" in style:
                chunks.append(f"## {text}")
            elif "heading 3" in style:
                chunks.append(f"### {text}")
            else:
                chunks.append(text)

        # Extract tables as simple text.
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    chunks.append(row_text)

        return "\n\n".join(chunks)[:max_chars]
    except ImportError:
        # python-docx not installed
        return ""
    except Exception:
        return ""


def extract_text_sample(path: Path, max_chars: int = 8000) -> str:
    try:
        size_bytes = path.stat().st_size
    except OSError:
        return ""
    if size_bytes > MAX_EXTRACT_BYTES:
        return f"[File too large for extraction: {size_bytes // 1024 // 1024}MB]"

    suffix = path.suffix.lower()
    if suffix == ".pptx" or suffix == ".ppt":
        return _read_pptx(path, max_chars)
    elif suffix == ".docx" or suffix == ".doc":
        return _read_docx(path, max_chars)
    elif suffix == ".pdf":
        return _read_pdf(path, max_chars)
    else:
        try:
            return path.read_text(encoding="utf-8", errors="ignore")[:max_chars]
        except OSError:
            return ""
